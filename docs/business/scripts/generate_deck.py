from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_pitch_deck():
    prs = Presentation()
    
    # Define corporate colors (Seed2Success Green)
    DARK_GREEN = RGBColor(11, 102, 35)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Seed2Success"
    subtitle.text = "The AI Operating System for Global Agriculture\nRaising $500,000 Seed"
    
    # Slide 2: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Problem: Farming is Broken"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "• Climate change is destroying crop yields."
    p = tf.add_paragraph()
    p.text = "• 500M Smallholder farmers rely on generational intuition."
    p = tf.add_paragraph()
    p.text = "• Existing AgTech is too expensive and requires hardware."
    
    # Slide 3: The Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Solution: A Multi-Modal AI Copilot"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "• Computer Vision: Detect crop diseases instantly."
    p = tf.add_paragraph()
    p.text = "• Machine Learning: Predict harvest yield via weather and soil."
    p = tf.add_paragraph()
    p.text = "• RAG LLM + Voice: Farmers 'talk' to expert AI in native languages."
    
    # Slide 4: The Moat
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Moat: Continuous Learning"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "• We use MLOps to capture actual harvest feedback."
    p = tf.add_paragraph()
    p.text = "• The model detects drift and retrains automatically."
    p = tf.add_paragraph()
    p.text = "• Microsoft FarmBeats does not have this proprietary data loop."

    # Slide 5: Business Model
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Business Model: B2B2C"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "• B2C (Free): Build the massive Data Flywheel."
    p = tf.add_paragraph()
    p.text = "• B2B ($10k-$50k/mo): License ML endpoints to Agri-Giants."
    p = tf.add_paragraph()
    p.text = "• B2G ($100k+/yr): Macro yield analytics for Governments."
    
    # Slide 6: Financials
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Financials & Ask"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "• Raising $500,000 Seed at $5M Post-Money Valuation."
    p = tf.add_paragraph()
    p.text = "• 28 Months of Runway ($17.5k monthly burn)."
    p = tf.add_paragraph()
    p.text = "• Break-Even projected at Month 22 via B2B APIs."
    
    prs.save("Seed2Success_Pitch_Deck.pptx")
    print("Successfully generated Seed2Success_Pitch_Deck.pptx")

if __name__ == "__main__":
    create_pitch_deck()
